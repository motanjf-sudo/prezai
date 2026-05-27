from fastapi import FastAPI, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
import anthropic
import requests
import json
import os
import io
import math
import base64
from PIL import Image, ImageDraw, ImageFont
import arabic_reshaper
from bidi.algorithm import get_display
from pptx import Presentation
from pptx.util import Inches

app = FastAPI()

class PresentationRequest(BaseModel):
    anthropic_key: str
    gemini_key: str = ""
    topic: str
    subtopic: str = ""
    audience: str = ""
    goal: str = ""
    style: str = "modern professional"
    tone: str = "حرفه‌ای"
    slides: int = 5
    color: str = "deep blue"
    theme: str = "dark"
    brand: str = ""
    hashtags: str = ""
    instagram: str = ""
    phone: str = ""
    website: str = ""

def fix_persian(text):
    try:
        reshaped = arabic_reshaper.reshape(str(text))
        return get_display(reshaped)
    except:
        return str(text)

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def get_font(size, bold=False):
    try:
        if bold:
            return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", size)
        return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size)
    except:
        return ImageFont.load_default()

def draw_rounded_rect(draw, x, y, w, h, radius, fill, border=None):
    draw.rounded_rectangle([x, y, x+w, y+h], radius=radius, fill=fill, outline=border, width=2)

def draw_arrow(draw, x1, y1, x2, y2, color, width=3):
    draw.line([(x1,y1),(x2,y2)], fill=color, width=width)
    angle = math.atan2(y2-y1, x2-x1)
    s = 12
    draw.polygon([
        (x2, y2),
        (int(x2-s*math.cos(angle-0.4)), int(y2-s*math.sin(angle-0.4))),
        (int(x2-s*math.cos(angle+0.4)), int(y2-s*math.sin(angle+0.4)))
    ], fill=color)

def ask_claude(client, prompt, max_tokens=1500):
    response = client.messages.create(
        model="claude-sonnet-4-5",
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}]
    )
    return response.content[0].text

def generate_image_gemini(gemini_key, prompt):
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-3.0-generate-002:predict?key={gemini_key}"
        body = {"instances": [{"prompt": prompt}], "parameters": {"sampleCount": 1, "aspectRatio": "16:9"}}
        response = requests.post(url, headers={"Content-Type": "application/json"}, json=body, timeout=60)
        if response.status_code == 200:
            data = response.json()
            img_data = base64.b64decode(data["predictions"][0]["bytesBase64Encoded"])
            return Image.open(io.BytesIO(img_data))
    except:
        pass
    return None

def render_slide(palette, elements, image=None, W=1280, H=720):
    bg_color = hex_to_rgb(palette["background"])
    img = Image.new("RGB", (W, H), bg_color)
    draw = ImageDraw.Draw(img)

    if image:
        image_resized = image.resize((W, H), Image.LANCZOS)
        img.paste(image_resized, (0, 0))
        overlay = Image.new("RGBA", (W, H), (0, 0, 0, 130))
        img = img.convert("RGBA")
        img = Image.alpha_composite(img, overlay)
        img = img.convert("RGB")
        draw = ImageDraw.Draw(img)

    for elem in elements:
        try:
            t = elem.get("type")
            x = int(elem.get("x", 0))
            y = int(elem.get("y", 0))
            w = int(elem.get("w", 100))
            h = int(elem.get("h", 50))
            color = hex_to_rgb(elem.get("color", palette["primary"]))
            radius = int(elem.get("radius", 10))

            if t == "header_band":
                draw.rectangle([0, y, W, y+h], fill=color)
            elif t == "footer_band":
                draw.rectangle([0, H-h, W, H], fill=color)
            elif t == "rect":
                border = hex_to_rgb(elem["border"]) if elem.get("border") else None
                draw_rounded_rect(draw, x, y, w, h, radius, color, border)
            elif t == "text":
                txt = fix_persian(elem.get("text", ""))
                size = int(elem.get("size", 18))
                font = get_font(size, elem.get("bold", False))
                draw.text((x, y), txt, fill=color, font=font)
            elif t == "text_wrapped":
                txt = fix_persian(elem.get("text", ""))
                size = int(elem.get("size", 16))
                max_w = int(elem.get("max_w", 300))
                font = get_font(size)
                words = txt.split()
                lines = []
                cur = ""
                for word in words:
                    test = cur + " " + word if cur else word
                    bbox = draw.textbbox((0,0), test, font=font)
                    if bbox[2]-bbox[0] <= max_w:
                        cur = test
                    else:
                        if cur: lines.append(cur)
                        cur = word
                if cur: lines.append(cur)
                for li, line in enumerate(lines):
                    draw.text((x, y+li*(size+4)), line, fill=color, font=font)
            elif t == "arrow":
                draw_arrow(draw, x, y, int(elem.get("x2",x+100)), int(elem.get("y2",y)), color, int(elem.get("line_width",3)))
            elif t == "divider":
                draw.rectangle([x, y, x+w, y+int(elem.get("thickness",3))], fill=color)
            elif t == "icon_box":
                bc = hex_to_rgb(elem.get("box_color", palette["primary"]))
                draw_rounded_rect(draw, x, y, w, h, radius, bc)
                tc = hex_to_rgb(elem.get("text_color", "#ffffff"))
                draw.text((x+10, y+10), fix_persian(elem.get("title","")), fill=tc, font=get_font(int(elem.get("title_size",16)), True))
                draw.text((x+10, y+36), fix_persian(elem.get("body","")), fill=tc, font=get_font(int(elem.get("body_size",13))))
            elif t == "number_highlight":
                bc = hex_to_rgb(elem.get("box_color", palette["accent"]))
                draw_rounded_rect(draw, x, y, w, h, radius, bc)
                draw.text((x+10, y+10), str(elem.get("number","")), fill=hex_to_rgb(elem.get("number_color","#ffffff")), font=get_font(int(elem.get("number_size",48)), True))
                draw.text((x+10, y+h-28), fix_persian(elem.get("label","")), fill=hex_to_rgb(elem.get("label_color","#ffffff")), font=get_font(int(elem.get("label_size",16))))
            elif t == "step_boxes":
                steps = elem.get("steps", [])
                sw = int(elem.get("step_w", 180))
                sh = int(elem.get("step_h", 100))
                gap = int(elem.get("gap", 40))
                colors = elem.get("colors", [palette["primary"], palette["accent"]])
                for si, step in enumerate(steps):
                    sx = x + si*(sw+gap)
                    sc = hex_to_rgb(colors[si % len(colors)])
                    draw_rounded_rect(draw, sx, y, sw, sh, 10, sc)
                    draw.text((sx+10, y+10), fix_persian(step.get("title","")), fill=(255,255,255), font=get_font(14, True))
                    draw.text((sx+10, y+34), fix_persian(step.get("body","")), fill=(220,220,220), font=get_font(12))
                    if si < len(steps)-1:
                        draw_arrow(draw, sx+sw+5, y+sh//2, sx+sw+gap-5, y+sh//2, (255,255,255), 2)
        except Exception as e:
            continue

    return img

def parse_json(text):
    start = text.find("{")
    end = text.rfind("}") + 1
    return json.loads(text[start:end])

@app.post("/generate")
async def generate_presentation(req: PresentationRequest):
    try:
        client = anthropic.Anthropic(api_key=req.anthropic_key.strip())

        structure_prompt = f"""You are a world-class presentation designer.
Create a stunning presentation structure.
Topic: {req.topic}
Subtopic: {req.subtopic}
Audience: {req.audience}
Goal: {req.goal}
Slides: {req.slides}
Color: {req.color}
Style: {req.style}
Theme: {req.theme}

Return ONLY this JSON:
{{
  "palette": {{"background":"#hex","primary":"#hex","secondary":"#hex","accent":"#hex","text":"#hex"}},
  "slides": [{{"number":1,"title":"Persian title","focus":"main point","image_prompt":"English image prompt, {req.style} style"}}]
}}"""

        structure_text = ask_claude(client, structure_prompt, 1000)
        structure = parse_json(structure_text)
        palette = structure["palette"]
        slides_info = structure["slides"]

        slide_images = []
        for slide_info in slides_info:
            elem_prompt = f"""Design ONE infographic slide. Canvas 1280x720px.
Title: {slide_info['title']}
Focus: {slide_info['focus']}
Style: {req.style}
Brand: {req.brand}
Hashtags: {req.hashtags}
Palette: bg={palette['background']}, primary={palette['primary']}, accent={palette['accent']}, text={palette['text']}

Use 8-10 elements. Return ONLY JSON: {{"elements":[...]}}

Types: header_band(y,h,color), footer_band(h,color), rect(x,y,w,h,color,radius), text(x,y,text,size,color,bold), text_wrapped(x,y,text,size,color,max_w), arrow(x,y,x2,y2,color,line_width), divider(x,y,w,color,thickness), icon_box(x,y,w,h,title,body,box_color,text_color), number_highlight(x,y,w,h,number,label,box_color,number_color,label_color), step_boxes(x,y,steps:[{{title,body}}],step_w,step_h,gap,colors:[])"""

            for attempt in range(3):
                try:
                    elem_text = ask_claude(client, elem_prompt, 2000)
                    elem_data = parse_json(elem_text)
                    elements = elem_data.get("elements", [])
                    break
                except:
                    elements = [
                        {"type": "header_band", "y": 0, "h": 80, "color": palette["primary"]},
                        {"type": "text", "x": 50, "y": 20, "text": slide_info["title"], "size": 36, "color": "#ffffff", "bold": True}
                    ]

            image = None
            if req.gemini_key and slide_info.get("image_prompt"):
                for attempt in range(2):
                    image = generate_image_gemini(req.gemini_key, slide_info["image_prompt"])
                    if image:
                        break

            img = render_slide(palette, elements, image)
            slide_images.append(img)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        for img in slide_images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            img_stream = io.BytesIO()
            img.save(img_stream, format="PNG")
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(13.33), Inches(7.5))

        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)

        from fastapi.responses import StreamingResponse
        return StreamingResponse(
            output_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=PrezAI.pptx"}
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

app.mount("/", StaticFiles(directory=".", html=True), name="static")
